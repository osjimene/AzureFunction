from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import date
from pptx.dml.color import RGBColor
import os

#Used to generate the powerpoint based on the ADO dataframe

def create_ppt(outfilename, data, RZTitle):
    dir = os.path.dirname(__file__)
    PPTXTemplate = os.path.join(dir, 'RedZoneInput.pptx')
    # outfilepath = os.environ["LocalTempFilePath"] + f" {outfilename}"
    outfilepath = os.path.join(dir, outfilename)
    exlaimIcon = os.path.join(dir, "exclaim.png")
    refreshIcon = os.path.join(dir, "refresh.png")
    resourceIcon = os.path.join(dir, "resource.png")
    
    """ Take the input powerpoint file and use it as the template for the output
    file.
    """
    prs = Presentation(PPTXTemplate)
    # Use the output from analyze_ppt to understand which layouts and placeholders
    # to use
    # Create a title slide first
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[12]
    title.text = "{} Monthly RedZone".format(RZTitle) #Need to fix this to be custom
    subtitle.text = "Date: {:%B %Y}".format(date.today()) 
        
    slidedata = data[['Issue','Req Date','Status','MSD ADO ID','MSD Owner','PG ADO','PG Owner','Comments']]

    #Numerate the amount of RedZone items to set the number of slides you will need. 
    rz_items = len(data)
    slide_number = (-(-rz_items//7))-1
    #Create New Slide using the Table Template
    slide_start = 0
    while slide_start <= slide_number: 
        redzone_item_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(redzone_item_layout)
        shapes = slide.shapes
        # shapes.title.text = 'RedZone Items'
        #Set the Table dimensions in the slide
        rows = cols = 8
        left = Inches(0.02)
        top = Inches(1.1)
        width = Inches(13.31)
        height = Inches(5.78)
        table = shapes.add_table(rows, cols, left, top, width, height).table
        #Set the Column widths
        table.columns[0].width= Inches(4.0)
        table.columns[1].width= Inches(1.02)
        table.columns[2].width= Inches(0.9)
        table.columns[3].width= Inches(0.89)
        table.columns[4].width= Inches(1.31)
        table.columns[5].width= Inches(0.8)
        table.columns[6].width= Inches(1.03)
        table.columns[7].width= Inches(3.38)
        table.rows[0].height = Inches(0.51)
        #Insert Top row values
        issue_title = table.cell(0,0)
        issue_req_date = table.cell(0,1)
        issue_status = table.cell(0,2)
        issue_id = table.cell(0,3)
        issue_owner = table.cell(0,4)
        issue_pg_id = table.cell(0,5)
        issue_pg_owner = table.cell(0,6)
        issue_comments = table.cell(0,7)
        #Adjust the text formatting on the top row
        #Add top row column values
        issue_title.text = "Issue"
        issue_req_date.text = "Req Date"
        issue_status.text = "Status"
        issue_id.text = "MSD ADO"
        issue_owner.text = "MSD Owner"
        issue_pg_id.text = "PG ADO"
        issue_pg_owner.text = "PG Owner"
        issue_comments.text = "Comments"
        #Adjust the Text formatting on the top row
        columns = [i for i in range(8)]
        for cols in columns:
            table.cell(0, cols).text_frame.paragraphs[0].font.size = Pt(12)
            table.cell(0, cols).text_frame.paragraphs[0].font.bold = True
            table.cell(0, cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
        rows = [x for x in range(1,8)]
       #Adjust the font size on the MSD Owner Column 
        for row in rows:
            table.cell(row, 4).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(row, 4).text_frame.paragraphs[0].font.bold = True
        #Adjust the font size on the PG Owner column
        for row in rows:
            table.cell(row, 6).text_frame.paragraphs[0].font.size = Pt(8)
            table.cell(row, 6).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
        slide_start += 1
        #Insert the moveable icons into the slides. 
        refresh = slide.shapes.add_picture(refreshIcon, left=Inches(7.71), top=Inches(7.21),height=Inches(0.17), width = Inches(0.23))
        exclaimation = slide.shapes.add_picture(exlaimIcon, left=Inches(6.51), top=Inches(7.16),height=Inches(0.28), width = Inches(0.05))
        resource = slide.shapes.add_picture(resourceIcon, left=Inches(5.18), top=Inches(7.18),height=Inches(0.27), width = Inches(0.26))
    #Insert values into the table.
    dataSlide = 1
    redzoneRow = 0
    while (redzoneRow < rz_items):
        populatedSlide = prs.slides[dataSlide]
        table = [shape for shape in populatedSlide.shapes if shape.has_table]
        for rows in range(1,8):
            for cols in range(0,8):
                #This adds hyperlinks to PG ADO Links.
                if 'https:' in str(slidedata.iloc[redzoneRow, cols]):
                    cell = table[0].table.cell(rows,cols)
                    cell_link = cell.text_frame.paragraphs[0].add_run()
                    cell_link.text = 'PG Link' 
                    hlink = cell_link.hyperlink
                    hlink.address = str(slidedata.iloc[redzoneRow, cols])
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(12)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Truncates the Character limit of the comments section to 200chars
                elif cols == 7:
                    condense = str(slidedata.iloc[redzoneRow, cols])[:220]
                    table[0].table.cell(rows,cols).text = condense 
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(9)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #This will add Hyperlinks to MSD ADO items.
                elif cols == 3:
                    cell = table[0].table.cell(rows,cols)
                    cell_link = cell.text_frame.paragraphs[0].add_run()
                    cell_link.text = str(slidedata.iloc[redzoneRow, cols]) 
                    hlink = cell_link.hyperlink
                    hlink.address = 'https://microsoftit.visualstudio.com/OneITVSO/_workitems/edit/'+ str(slidedata.iloc[redzoneRow, cols])
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(12)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'               
                #Logig for the status cell to be formatted properly.
                elif "RZ-" in str(slidedata.iloc[redzoneRow, cols]):
                    if 'Red' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0xFF, 0x4C, 0x4C)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    elif 'Yellow' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0xFF, 0xD3, 0x4C)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    elif 'Green' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0xB2, 0xDE, 0x84)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    elif 'Blue' in str(slidedata.iloc[redzoneRow, cols]):
                        table[0].table.cell(rows,cols).fill.solid()
                        table[0].table.cell(rows,cols).fill.fore_color.rgb = RGBColor(0x4C, 0xC8, 0xF4)
                        table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(11)
                        table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                    else:
                        continue
                else:
                    table[0].table.cell(rows,cols).text = str(slidedata.iloc[redzoneRow, cols]) 
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.size = Pt(12)
                    table[0].table.cell(rows,cols).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Adjust the font size on the MSD Owner Column 
                table[0].table.cell(rows, 4).text_frame.paragraphs[0].font.size = Pt(8)
                table[0].table.cell(rows, 4).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Adjust the font size on the PGADO Column 
                table[0].table.cell(rows, 5).text_frame.paragraphs[0].font.size = Pt(8)
                table[0].table.cell(rows, 5).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
                #Adjust the font size on the PG Owner column         
                table[0].table.cell(rows, 6).text_frame.paragraphs[0].font.size = Pt(8)
                table[0].table.cell(rows, 6).text_frame.paragraphs[0].font.name = 'Segoe UI (Body)'
            redzoneRow += 1
            if redzoneRow >= rz_items:
                break
            else:
                continue
        dataSlide += 1


    prs.save(outfilepath)
    return outfilepath